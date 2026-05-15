# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Cm, Pt


SOURCE_NAME = "골든듀 사용자 메뉴얼.pptx"
OUTPUT = Path(__file__).resolve().parents[1] / "goldendew_coupon_table_added.pptx"
TABLE_SHAPE_NAME = "CouponConditionsTable"


def find_source() -> Path:
    base = Path.home() / "Documents"
    for path in base.rglob(SOURCE_NAME):
        if not path.name.startswith("~$"):
            return path
    raise FileNotFoundError(SOURCE_NAME)


def remove_shape_by_name(slide, name: str) -> None:
    for shape in list(slide.shapes):
        if shape.name == name:
            slide.shapes._spTree.remove(shape._element)


def set_cell(cell, text, fill, font_color, font_size=8.4, bold=False, align=PP_ALIGN.LEFT):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Cm(0.08)
    cell.margin_right = Cm(0.08)
    cell.margin_top = Cm(0.04)
    cell.margin_bottom = Cm(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    text_frame = cell.text_frame
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = "맑은 고딕"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color


def set_cell_border(cell, color: RGBColor, width="9525") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        for old in tc_pr.findall(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"):
            tc_pr.remove(old)
        line = parse_xml(
            f'<a:{tag} {nsdecls("a")} w="{width}">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/>'
            f'</a:{tag}>'
        )
        tc_pr.append(line)


def add_coupon_table(prs: Presentation) -> None:
    slide = prs.slides[22]
    remove_shape_by_name(slide, TABLE_SHAPE_NAME)

    rows = [
        ["항목", "의미", "설정 방식", "적용 기준 / 예시"],
        [
            "사용 가능 주문 구분",
            "쿠폰을 적용할 수 있는 주문 유형을 제한합니다.",
            "적용할 주문 조건을 체크합니다. 아무것도 선택하지 않으면 전체 주문 유형에 적용 가능합니다.",
            "고재, 진열, 수선, 가주문, B2B 중 필요한 항목만 선택",
        ],
        [
            "적용 가능 최대 할인율",
            "판매 금액 대비 기존 할인율 조건입니다.",
            "설정한 할인율 이상 할인된 판매 건에서 쿠폰을 적용할 수 있습니다.",
            "30%로 설정 시 판매 금액의 30% 이상 할인이 발생한 경우 적용 가능",
        ],
        [
            "사용 가능 최소 금액",
            "판매 금액 기준의 최소/최대 금액 조건입니다.",
            "비교 조건(이하/이상)에 따라 입력 금액과 판매 금액이 조건에 맞아야 합니다.",
            "100,000원 이상 설정 시 판매 금액이 100,000원 이상일 때 적용 가능",
        ],
        [
            "사용 가능 최소 수량",
            "판매 개수 기준의 최소/최대 수량 조건입니다.",
            "비교 조건(이하/이상)에 따라 구매 수량이 조건에 맞아야 합니다.",
            "3개 이상 설정 시 3개 이상 구매한 경우 적용 가능",
        ],
    ]

    left = Cm(0.98)
    top = Cm(4.45)
    width = Cm(31.0)
    height = Cm(10.7)

    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table_shape.name = TABLE_SHAPE_NAME
    table = table_shape.table

    col_widths = [Cm(4.5), Cm(7.2), Cm(9.4), Cm(9.9)]
    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = col_width

    table.rows[0].height = Cm(0.95)
    for idx in range(1, len(rows)):
        table.rows[idx].height = Cm(2.35)

    purple = RGBColor(124, 72, 132)
    light_purple = RGBColor(246, 240, 248)
    pale_purple = RGBColor(251, 248, 252)
    white = RGBColor(255, 255, 255)
    dark_text = RGBColor(45, 42, 48)
    border = RGBColor(201, 169, 208)

    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, text, purple, white, font_size=9.0, bold=True, align=PP_ALIGN.CENTER)
            elif c == 0:
                set_cell(cell, text, light_purple, purple, font_size=8.4, bold=True, align=PP_ALIGN.CENTER)
            else:
                fill = pale_purple if r % 2 == 0 else white
                set_cell(cell, text, fill, dark_text, font_size=8.2)

            set_cell_border(cell, border)


def main() -> None:
    source = find_source()
    prs = Presentation(str(source))
    add_coupon_table(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(str(source))
    print(str(OUTPUT))


if __name__ == "__main__":
    main()
